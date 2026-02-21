#!/usr/bin/env python3
import os
import re
from pathlib import Path
from typing import List, Dict

import pdfplumber
import docx
from pptx import Presentation

# =========================
# DOCUMENT PROCESSOR
# =========================

class DocumentProcessor:
    def extract_text(self, file_path: str) -> str:
        if not os.path.exists(file_path):
            raise FileNotFoundError(file_path)

        ext = Path(file_path).suffix.lower()

        if ext == ".pdf":
            return self._from_pdf(file_path)
        elif ext == ".docx":
            return self._from_docx(file_path)
        elif ext == ".pptx":
            return self._from_pptx(file_path)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported format: {ext}")

    def _from_pdf(self, file_path):
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"
        return text

    def _from_docx(self, file_path):
        doc = docx.Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text


# =========================
# CHAPTER SEGMENTER
# =========================

class ChapterSegmenter:
    CHAPTER_PATTERN = r"(chapter\s+\d+|unit\s+\d+|week\s+\d+|section\s+\d+)"

    def split_into_chapters(self, text: str) -> List[str]:
        parts = re.split(self.CHAPTER_PATTERN, text, flags=re.IGNORECASE)

        chapters = []
        current = ""

        for part in parts:
            if re.match(self.CHAPTER_PATTERN, part, re.IGNORECASE):
                if current:
                    chapters.append(current.strip())
                current = part
            else:
                current += " " + part

        if current.strip():
            chapters.append(current.strip())

        return [c for c in chapters if len(c.split()) > 40]


# =========================
# EXPLANATION ENGINE
# =========================

class LectureExplainer:

    def generate_explanation(self, text: str, mode: str) -> str:
        sentences = re.split(r'[.!?]', text)
        key_points = " ".join(sentences[:5])

        if mode == "simple":
            return (
                "This chapter introduces the main ideas in a simple way. "
                f"The key focus is: {key_points}. "
                "Pay attention to the basic concepts and definitions."
            )

        elif mode == "lecture":
            return (
                "This chapter presents important academic concepts. "
                f"It discusses: {key_points}. "
                "These ideas form the foundation for deeper understanding of the subject."
            )

        elif mode == "exam":
            return (
                "For examination purposes, focus on the core concepts discussed. "
                f"The most important ideas include: {key_points}. "
                "Understand definitions, relationships, and applications."
            )

        elif mode == "detailed":
            return (
                "This chapter provides a comprehensive explanation of the topic. "
                f"It explores: {key_points}. "
                "The discussion connects theory with practical understanding."
            )

        else:
            return "This chapter explains key concepts and their applications."


# =========================
# AUDIO GENERATOR
# =========================

class AudioGenerator:

    def __init__(self, voice_type="default", speech_rate="normal"):
        self.voice_type = voice_type
        self.speech_rate = speech_rate

    def text_to_audio(self, text: str, output_file: str):
        try:
            if self.voice_type == "african_male_online":
                self._edge_tts(text, output_file)
            else:
                self._offline_tts(text, output_file)
        except Exception:
            self._offline_tts(text, output_file)

    def _offline_tts(self, text, output_file):
        import pyttsx3

        engine = pyttsx3.init()

        for voice in engine.getProperty("voices"):
            if "male" in voice.name.lower():
                engine.setProperty("voice", voice.id)

        if self.speech_rate == "slow":
            engine.setProperty("rate", 140)

        engine.save_to_file(text, output_file)
        engine.runAndWait()

    def _edge_tts(self, text, output_file):
        import edge_tts
        import asyncio

        voice = "en-NG-EzinneNeural"

        async def run():
            communicate = edge_tts.Communicate(text, voice)
            await communicate.save(output_file)

        asyncio.run(run())

    def combine(self, files: List[str], output_file: str):
        from pydub import AudioSegment

        combined = AudioSegment.empty()
        for f in files:
            combined += AudioSegment.from_file(f)
            combined += AudioSegment.silent(duration=500)

        combined.export(output_file, format="mp3")


# =========================
# MAIN PROCESSOR
# =========================

class CourseAudioProcessor:

    def __init__(self):
        self.doc = DocumentProcessor()
        self.segmenter = ChapterSegmenter()
        self.explainer = LectureExplainer()

    def process_document(
        self,
        file_path: str,
        output_dir: str = "./output",
        explanation_mode: str = "lecture",
        voice_type: str = "default",
        speech_rate: str = "normal"
    ) -> Dict:

        os.makedirs(output_dir, exist_ok=True)

        text = self.doc.extract_text(file_path)
        chapters = self.segmenter.split_into_chapters(text)

        if not chapters:
            raise ValueError("No chapters detected")

        audio = AudioGenerator(voice_type, speech_rate)

        temp_files = []

        for i, chapter in enumerate(chapters, 1):
            read_file = os.path.join(output_dir, f"chapter_{i}_read.mp3")
            explain_file = os.path.join(output_dir, f"chapter_{i}_explain.mp3")

            audio.text_to_audio(f"Chapter {i}. {chapter}", read_file)

            explanation = self.explainer.generate_explanation(
                chapter,
                explanation_mode
            )
            audio.text_to_audio(f"Explanation. {explanation}", explain_file)

            temp_files.extend([read_file, explain_file])

        final_audio = os.path.join(output_dir, "course_audio.mp3")
        audio.combine(temp_files, final_audio)

        return {
            "success": True,
            "chapters_processed": len(chapters),
            "audio_file": final_audio
        }


